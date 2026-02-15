# streamlit_app.py
# -*- coding: utf-8 -*-
import os, re, json, socket, io, base64, urllib.parse, traceback, logging
from pathlib import Path
from typing import Optional, Tuple, Dict, Any, List
from urllib.parse import urlparse, urlunparse
from copy import deepcopy

import requests
import streamlit as st

# ====== PPT / Images / PDF ======
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Emu
from PIL import Image, ImageChops
import pypdfium2 as pdfium  # PDF -> images

from google.auth.exceptions import RefreshError

# ====== (Optionnel) Gmail OAuth ‚Äì pour r√©cup√©rer le PDF Looker par email ======
try:
    from google.oauth2.credentials import Credentials
    from googleapiclient.discovery import build
    from googleapiclient.errors import HttpError
    GMAIL_AVAILABLE = True
except Exception:
    GMAIL_AVAILABLE = False


# =====================================
# Page config
# =====================================
st.set_page_config(
    page_title="Export Pr√©sentation (Tableau + Looker)",
    page_icon="üìë",
    layout="wide"
)

APP_DIR = Path(__file__).parent.resolve()

# ----- Dossiers techniques -----
DATA_DIR = APP_DIR / "data"
DATA_DIR.mkdir(parents=True, exist_ok=True)

CAPTURES_DIR = APP_DIR / "captures"
CONFIG_DIR = APP_DIR / "config"
LOG_DIR = APP_DIR / "log"

for d in (CAPTURES_DIR, CONFIG_DIR, LOG_DIR):
    d.mkdir(parents=True, exist_ok=True)

# Cache local (sur cet appareil) ‚Äì on le met dans config/
LOCAL_CREDS_PATH = CONFIG_DIR / "local_device_creds.json"

# Fichier de log debug
DEBUG_LOG_PATH = LOG_DIR / "debug.log"
logging.basicConfig(
    filename=str(DEBUG_LOG_PATH),
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)
logger.info("Application Streamlit d√©marr√©e")

# PIN (optionnel)
APP_PIN = (os.getenv("APP_PIN") or "").strip()

# Version API Tableau
TABLEAU_API_VERSION = os.getenv("TABLEAU_API_VERSION", "3.22")

# ===== Placeholders (m√™mes noms que ton script valid√©) =====
PH_TBL, PH_TBL_TITLE = "PH_TBL", "PH_TBL_TITLE"
PH_LKR_BASE, PH_LKR_TITLE_BASE = "PH_LKR", "PH_LKR_TITLE"
PH_IMAGE_MAIN, PH_TITLE_MAIN = "PH_IMAGE_MAIN", "PH_TITLE"

# Rognage top : pas de slider, contr√¥lable par variable d'env
TOPBAR_CROP_PCT = float(os.getenv("TOPBAR_CROP_PCT", "0.10"))  # 12% par d√©faut


# --- Helper pour lire les secrets depuis ENV ou st.secrets ---
def _get_secret(name: str, default: str = "") -> str:
    """
    Essaie d'abord os.getenv, sinon st.secrets[name] si pr√©sent, sinon renvoie default.
    Toujours une string.
    """
    try:
        v = os.getenv(name)
        if v:
            return str(v)
        if hasattr(st, "secrets"):
            v = st.secrets.get(name)
            if v:
                return str(v)
    except Exception:
        pass
    return default


# Ajustement d'image dans le placeholder :
# - "contain" : pas de crop (centr√© avec √©ventuelles marges)
# - "cover" : recouvre la zone (peut rogner)
FIT_MODE = (os.getenv("IMAGE_FIT_MODE", "contain") or "contain").lower()

# --- Looker: ajustements sp√©cifiques ---
# Pas de crop lat√©ral, pas de trim par d√©faut: on travaille en 'contain'
LKR_FIT_MODE = (os.getenv("LKR_FIT_MODE", "contain") or "contain").lower()
LKR_TRIM = (os.getenv("LKR_TRIM", "false").lower() == "true")  # false par d√©faut

# Rognage par pourcentage du PDF Looker (utile pour couper une ligne d'ent√™te/pied)
LKR_CROP_TOP = float(os.getenv("LKR_CROP_TOP", "0"))   # ex: 0.08 (= 8 %)
LKR_CROP_BOTTOM = float(os.getenv("LKR_CROP_BOTTOM", "0"))
LKR_CROP_LEFT = float(os.getenv("LKR_CROP_LEFT", "0.08"))
LKR_CROP_RIGHT = float(os.getenv("LKR_CROP_RIGHT", "0.08"))

# ====== Gmail ENV (si tu l'actives) ======
GMAIL_CLIENT_ID = _get_secret("GMAIL_CLIENT_ID")
GMAIL_CLIENT_SECRET = _get_secret("GMAIL_CLIENT_SECRET")
GMAIL_REFRESH_TOKEN = _get_secret("GMAIL_REFRESH_TOKEN")
GMAIL_USER = _get_secret("GMAIL_USER", "me") or "me"
DEFAULT_SENDER = "looker-studio-noreply@google.com"


# =====================================
# Helpers cache local
# =====================================
def _read_local_device_cache() -> dict:
    try:
        if LOCAL_CREDS_PATH.exists():
            return json.loads(LOCAL_CREDS_PATH.read_text(encoding="utf-8"))
    except Exception as e:
        logger.warning("Erreur lecture cache local_device_creds.json : %s", e)
    return {}


def _write_local_device_cache(d: dict) -> None:
    try:
        LOCAL_CREDS_PATH.write_text(
            json.dumps(d, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
    except Exception as e:
        logger.error("Erreur √©criture cache local_device_creds.json : %s", e)


# ---- Tableau creds ----
def load_local_device_creds() -> Optional[dict]:
    cache = _read_local_device_cache()
    return cache.get("tableau")


def save_local_device_creds(payload: dict) -> None:
    cache = _read_local_device_cache()
    cache["tableau"] = payload
    _write_local_device_cache(cache)


def clear_local_device_creds():
    cache = _read_local_device_cache()
    if "tableau" in cache:
        cache.pop("tableau")
    _write_local_device_cache(cache)


# ---- Looker sources (multi-sources Gmail/URL) ----
def load_looker_sources() -> List[dict]:
    cache = _read_local_device_cache()
    return cache.get("looker_sources", [])


def save_looker_sources(sources: List[dict]) -> None:
    cache = _read_local_device_cache()
    cache["looker_sources"] = sources
    _write_local_device_cache(cache)


def clear_looker_sources():
    cache = _read_local_device_cache()
    if "looker_sources" in cache:
        cache.pop("looker_sources")
    _write_local_device_cache(cache)


# =====================================
# Normalisation & validation host
# =====================================
def _normalize_server(url: str) -> str:
    url = (url or "").strip()
    if not url:
        return url
    if not re.match(r"^https?://", url, flags=re.I):
        url = "https://" + url
    url = url.rstrip("/")
    return url


def _autocorrect_tableau_pod_hostname(server_url: str) -> Tuple[str, Optional[str]]:
    """
    Corrige 'eu-west-la' -> 'eu-west-1a' si besoin.
    """
    try:
        parsed = urlparse(server_url)
        host = (parsed.hostname or "").lower()
        fixed_host = host
        if re.search(r"\beu-west-la\.online\.tableau\.com$", host):
            fixed_host = host.replace("eu-west-la.", "eu-west-1a.")
        if fixed_host != host:
            corrected = urlunparse(
                (parsed.scheme, fixed_host, parsed.path.rstrip("/"), "", "", "")
            )
            return corrected, f"Host corrig√© automatiquement : {host} ‚Üí {fixed_host}"
    except Exception:
        pass
    return server_url, None


def _check_dns(host: str) -> Tuple[bool, Optional[str]]:
    try:
        socket.getaddrinfo(host, 443)
        return True, None
    except Exception as e:
        return False, f"√âchec de r√©solution DNS pour '{host}': {e}"


# =====================================
# UI identifiants (Tableau)
# =====================================
def tableau_credentials_ui(
    saved: dict | None = None, key_prefix: str = "tbl_main"
) -> Tuple[str, str, str, str]:
    saved = saved or {}

    def k(name: str) -> str:
        return f"cred_{key_prefix}_{name}"

    c1, c2, c3, c4 = st.columns(4)
    server_in = c1.text_input(
        "Server",
        value=saved.get("server", "https://eu-west-1a.online.tableau.com"),
        key=k("server"),
    )
    site = c2.text_input(
        "Site (contentUrl)",
        value=saved.get("site", "autobiz"),
        key=k("site"),
    )
    pat_name = c3.text_input(
        "PAT Name",
        value=saved.get("pat_name", ""),
        key=k("pat_name"),
    )
    has_saved_secret = bool(saved.get("pat_secret"))
    pat_secret_input = c4.text_input(
        "PAT Secret",
        value="",
        type="password",
        key=k("pat_secret"),
        placeholder=("‚Ä¢‚Ä¢‚Ä¢‚Ä¢‚Ä¢‚Ä¢ (d√©j√† enregistr√©)" if has_saved_secret else ""),
    )

    server = _normalize_server(server_in)
    server, info = _autocorrect_tableau_pod_hostname(server)
    if info:
        st.info(info)

    pat_secret = pat_secret_input if pat_secret_input else saved.get("pat_secret", "")

    if server:
        host = urlparse(server).hostname or ""
        ok_dns, msg_dns = _check_dns(host)
        st.caption(f"‚úÖ DNS OK pour {host}" if ok_dns else f"‚ö†Ô∏è {msg_dns}")

    return server.strip(), site.strip(), pat_name.strip(), pat_secret.strip()


# =====================================
# Acc√®s (PIN optionnel)
# =====================================
def guard_access() -> bool:
    if not APP_PIN:
        return True

    st.sidebar.header("Acc√®s")
    pin = st.sidebar.text_input("Code d'acc√®s (PIN)", type="password")
    if st.sidebar.button("Entrer"):
        if pin == APP_PIN:
            st.session_state["pin_ok"] = True
            st.rerun()
        else:
            st.sidebar.error("Code incorrect.")
    return bool(st.session_state.get("pin_ok"))


# =====================================
# API Tableau (PAT)
# =====================================
JSON_HEADERS = {"Accept": "application/json", "Content-Type": "application/json"}


def _json_or_raise(resp: requests.Response) -> Dict[str, Any]:
    ct = (resp.headers.get("Content-Type") or "").lower()
    if "application/json" in ct:
        return resp.json()
    text = resp.text.strip()
    snippet = (text[:800] + "‚Ä¶") if len(text) > 800 else text
    resp.raise_for_status()
    raise RuntimeError(f"R√©ponse non-JSON (Content-Type={ct}). Extrait: {snippet}")


class TableauSession:
    def __init__(
        self, server: str, site: str, pat_name: str, pat_secret: str, api_version: str
    ):
        self.server = server.rstrip("/")
        self.site_content_url = site or ""
        self.pat_name = pat_name
        self.pat_secret = pat_secret
        self.api_version = api_version
        self.token = None
        self.site_id = None
        self.user_id = None

    @property
    def base(self) -> str:
        return f"{self.server}/api/{self.api_version}"

    def signin(self):
        logger.info(
            "Tentative de connexion Tableau : server=%s, site=%s, pat_name=%s",
            self.server,
            self.site_content_url,
            self.pat_name,
        )
        r = requests.post(
            f"{self.base}/auth/signin",
            headers=JSON_HEADERS,
            json={
                "credentials": {
                    "personalAccessTokenName": self.pat_name,
                    "personalAccessTokenSecret": self.pat_secret,
                    "site": {"contentUrl": self.site_content_url},
                }
            },
            timeout=30,
        )
        if r.status_code in (401, 403):
            logger.error(
                "Authentification Tableau refus√©e (status=%s)", r.status_code
            )
            raise RuntimeError("Authentification refus√©e (401/403).")
        c = _json_or_raise(r).get("credentials", {})
        self.token = c.get("token")
        self.site_id = (c.get("site") or {}).get("id")
        self.user_id = (c.get("user") or {}).get("id")
        if not (self.token and self.site_id and self.user_id):
            logger.error("Connexion Tableau incompl√®te : token/site_id/user_id manquants")
            raise RuntimeError("Connexion Tableau incompl√®te.")
        logger.info(
            "Connexion Tableau OK : site_id=%s, user_id=%s",
            self.site_id,
            self.user_id,
        )

    def signout(self):
        if not self.token:
            return
        try:
            requests.post(
                f"{self.base}/auth/signout",
                headers={"X-Tableau-Auth": self.token, **JSON_HEADERS},
                timeout=15,
            )
        finally:
            self.token = None

    def _headers(self) -> Dict[str, str]:
        if not self.token:
            raise RuntimeError("Non connect√©.")
        return {"X-Tableau-Auth": self.token, **JSON_HEADERS}

    def list_workbooks_for_user(self, page_size=1000) -> List[Dict[str, Any]]:
        r = requests.get(
            f"{self.base}/sites/{self.site_id}/users/{self.user_id}/workbooks",
            headers=self._headers(),
            params={"pageSize": page_size},
            timeout=30,
        )
        if r.status_code in (401, 403):
            raise RuntimeError("Acc√®s refus√© workbooks.")
        wbs = _json_or_raise(r).get("workbooks", {}).get("workbook", []) or []
        out = [
            {
                "id": wb.get("id"),
                "name": wb.get("name"),
                "projectName": (wb.get("project", {}) or {}).get("name")
                or wb.get("projectName"),
            }
            for wb in wbs
        ]
        out.sort(key=lambda x: (x.get("name") or "").lower())
        return out

    def list_views_for_workbook(
        self, workbook_id: str, page_size=1000
    ) -> List[Dict[str, Any]]:
        r = requests.get(
            f"{self.base}/sites/{self.site_id}/workbooks/{workbook_id}/views",
            headers=self._headers(),
            params={"pageSize": page_size},
            timeout=30,
        )
        if r.status_code in (401, 403):
            raise RuntimeError("Acc√®s refus√© vues.")
        vs = _json_or_raise(r).get("views", {}).get("view", []) or []
        out = [
            {"id": v.get("id"), "name": v.get("name"), "contentUrl": v.get("contentUrl")}
            for v in vs
        ]
        out.sort(key=lambda x: (x.get("name") or "").lower())
        return out

    # --- Exports PPT ---
    def export_view_ppt(self, view_id: str) -> Optional[bytes]:
        return self._download_binary(
            f"{self.base}/sites/{self.site_id}/views/{view_id}/powerpoint"
        )

    def export_workbook_ppt(self, wb_id: str) -> Optional[bytes]:
        return self._download_binary(
            f"{self.base}/sites/{self.site_id}/workbooks/{wb_id}/powerpoint"
        )

    def _download_binary(self, url: str) -> Optional[bytes]:
        accepts = [
            "application/octet-stream",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "*/*",
        ]
        for a in accepts:
            try:
                r = requests.get(
                    url,
                    headers={"Accept": a, **self._headers()},
                    stream=True,
                    timeout=300,
                )
                if r.status_code == 200:
                    buf = io.BytesIO()
                    for ch in r.iter_content(8192):
                        if ch:
                            buf.write(ch)
                    return buf.getvalue()
                if r.status_code in (400, 406, 415):
                    continue
                r.raise_for_status()
            except Exception as e:
                logger.warning(
                    "Erreur download_binary (Accept=%s) : %s",
                    a,
                    e,
                )
                continue

        # tentative sans header Accept
        try:
            r = requests.get(
                url,
                headers=self._headers(),
                stream=True,
                timeout=300,
            )
            if r.status_code == 200:
                buf = io.BytesIO()
                for ch in r.iter_content(8192):
                    if ch:
                        buf.write(ch)
                return buf.getvalue()
        except Exception as e:
            logger.error("Erreur download_binary sans Accept : %s", e)
        return None


# --- Cache Streamlit (catalogue) ---
@st.cache_data(show_spinner=False, ttl=600)
def cached_fetch_workbooks(
    server: str, site: str, pat_name: str, pat_secret: str, api_version: str
) -> List[Dict[str, Any]]:
    sess = TableauSession(server, site, pat_name, pat_secret, api_version)
    try:
        sess.signin()
        return sess.list_workbooks_for_user()
    finally:
        sess.signout()


@st.cache_data(show_spinner=False, ttl=600)
def cached_fetch_views(
    server: str,
    site: str,
    pat_name: str,
    pat_secret: str,
    api_version: str,
    workbook_id: str,
) -> List[Dict[str, Any]]:
    sess = TableauSession(server, site, pat_name, pat_secret, api_version)
    try:
        sess.signin()
        return sess.list_views_for_workbook(workbook_id)
    finally:
        sess.signout()


# =====================================
# Looker ‚Äî UI multi-sources (Gmail / URL) + filtres
# =====================================
def build_gmail_query(
    senders_csv: str | None, subjects_csv: str | None, label: str | None, days: int
) -> str:
    base = ["has:attachment", "filename:pdf", f"newer_than:{int(days)}d"]
    if label and label.strip():
        base.append(label if label.lower().startswith("label:") else f"label:{label.strip()}")

    use_senders = (senders_csv or "").strip() or DEFAULT_SENDER
    senders = [s.strip() for s in (use_senders or "").split(",") if s.strip()]
    subjects = [s.strip() for s in (subjects_csv or "").split(",") if s.strip()]

    ors = []
    if senders:
        ors.append("(" + " OR ".join([f"from:{s}" for s in senders]) + ")")
    if subjects:
        ors.append("(" + " OR ".join([f'subject:"{s}"' for s in subjects]) + ")")

    return " ".join(base + ors)


def looker_sources_ui():
    st.subheader("üìß Sources Looker")

    if "lkr_sources" not in st.session_state:
        st.session_state["lkr_sources"] = load_looker_sources() or [{"mode": "gmail"}]

    def draw_source(idx: int, src: dict):
        st.markdown(f"**Rapport {idx+1}**")
        cols = st.columns([1, 1, 1, 1, 1])

        mode = cols[0].selectbox(
            "Mode",
            ["gmail", "url"],
            index=(0 if src.get("mode", "gmail") == "gmail" else 1),
            key=f"lkr_mode_{idx}",
        )
        src["mode"] = mode

        if mode == "gmail":
            src["senders"] = cols[1].text_input(
                "Exp√©diteur(s)",
                value=src.get("senders", ""),
                key=f"lkr_send_{idx}",
            )
            src["subjects"] = cols[2].text_input(
                "Sujet(s)",
                value=src.get("subjects", "Looker,Data Studio,Rapport"),
                key=f"lkr_subj_{idx}",
            )
            src["label"] = cols[3].text_input(
                "Label",
                value=src.get("label", ""),
                key=f"lkr_lab_{idx}",
            )
            src["days"] = cols[4].number_input(
                "Jours",
                1,
                365,
                int(src.get("days", 60)),
                key=f"lkr_days_{idx}",
            )
            q = build_gmail_query(
                src.get("senders", ""),
                src.get("subjects", ""),
                src.get("label", ""),
                int(src.get("days", 60)),
            )
            st.caption(f"üîé Requ√™te Gmail g√©n√©r√©e : {q}")

        elif mode == "url":
            src["url"] = cols[1].text_input(
                "URL PDF publique",
                value=src.get("url", ""),
                key=f"lkr_url_{idx}",
                placeholder="https://‚Ä¶/export.pdf",
            )

        col_del = st.columns([1, 5, 1])
        if col_del[2].button("üóëÔ∏è Supprimer", key=f"del_src_{idx}"):
            st.session_state["lkr_sources"].pop(idx)
            st.rerun()

    for i, src in enumerate(st.session_state["lkr_sources"]):
        with st.container(border=True):
            draw_source(i, src)

    col_actions = st.columns([1, 1, 1])
    if col_actions[0].button("‚ûï Ajouter un rapport"):
        st.session_state["lkr_sources"].append({"mode": "gmail"})
        st.rerun()
    if col_actions[1].button("üíæ Enregistrer les sources"):
        save_looker_sources(st.session_state["lkr_sources"])
        st.success("Sources Looker enregistr√©es localement.")
    if col_actions[2].button("üóëÔ∏è R√©initialiser les sources"):
        clear_looker_sources()
        st.session_state["lkr_sources"] = [{"mode": "gmail"}]
        st.success("Sources Looker r√©initialis√©es.")
        st.rerun()


# =====================================
# Templates : d√©couverte & lecture
# =====================================
def discover_templates() -> Dict[str, Path]:
    candidates = [
        APP_DIR,
        APP_DIR / "templates",
        APP_DIR / "Files",
        APP_DIR / "files",
        Path("/app"),
        Path("/tmp"),
    ]
    out: Dict[str, Path] = {}
    for base in candidates:
        if not base.exists():
            continue
        for p in base.rglob("*.pptx"):
            if p.is_file() and not p.name.startswith("export_"):
                label = (
                    f"{p.name} ({str(p.relative_to(APP_DIR)) if APP_DIR in p.parents else str(p)})"
                )
                out[label] = p
    return out


def get_template_bytes(choice: str, discovered: Dict[str, Path]):
    if choice and choice in discovered:
        p = discovered[choice]
        try:
            return p.read_bytes(), choice
        except Exception:
            st.error("Lecture du template impossible.")
            st.code(traceback.format_exc())
    return None, None


# =====================================
# PDF / Images utils + remplissage PPT
# =====================================
def pdf_to_png_bytes(pdf_bytes: bytes, zoom: float = 2.0) -> List[bytes]:
    doc = pdfium.PdfDocument(io.BytesIO(pdf_bytes))
    images: List[bytes] = []
    try:
        for i in range(len(doc)):
            page = doc[i]
            pil = page.render(scale=zoom).to_pil()
            buf = io.BytesIO()
            pil.save(buf, format="PNG")
            images.append(buf.getvalue())
            page.close()
    finally:
        doc.close()
    return images


def _bytes_to_pil(b: bytes) -> Image.Image:
    return Image.open(io.BytesIO(b)).convert("RGB")


def _pil_to_bytes(img: Image.Image, fmt="PNG") -> bytes:
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    return buf.getvalue()


def trim_whitespace(pil_img: Image.Image) -> Image.Image:
    bg = Image.new(pil_img.mode, pil_img.size, (255, 255, 255))
    diff = ImageChops.difference(pil_img, bg)
    bbox = diff.getbbox()
    return pil_img.crop(bbox) if bbox else pil_img


def crop_top_bar(pil_img: Image.Image, topbar_pct: float = 0.0) -> Image.Image:
    topbar_pct = max(0.0, min(0.5, float(topbar_pct or 0.0)))
    if topbar_pct <= 0:
        return pil_img
    w, h = pil_img.size
    cut = int(h * topbar_pct)
    return pil_img.crop((0, cut, w, h)) if cut > 0 else pil_img


def add_picture_fit(
    slide,
    img_bytes: bytes,
    left,
    top,
    width,
    height,
    *,
    trim=True,
    topbar_crop_pct=0.0,
    fit="contain",
    pre_crop_pct=None,
):
    pil = _bytes_to_pil(img_bytes)

    # Recadrage doux par pourcentage AVANT tout (utile pour Looker)
    if pre_crop_pct:
        top_pct = max(0.0, min(0.49, float(pre_crop_pct.get("top", 0.0))))
        bottom_pct = max(0.0, min(0.49, float(pre_crop_pct.get("bottom", 0.0))))
        left_pct = max(0.0, min(0.49, float(pre_crop_pct.get("left", 0.0))))
        right_pct = max(0.0, min(0.49, float(pre_crop_pct.get("right", 0.0))))
        w, h = pil.size
        x1 = int(w * left_pct)
        y1 = int(h * top_pct)
        x2 = int(w * (1.0 - right_pct))
        y2 = int(h * (1.0 - bottom_pct))
        if x2 > x1 and y2 > y1:
            pil = pil.crop((x1, y1, x2, y2))

    # Trim + rognage haut (r√©serv√©s √† Tableau ; pour Looker on les d√©sactive via param√®tres)
    if trim:
        pil = trim_whitespace(pil)
    if topbar_crop_pct and topbar_crop_pct > 0:
        pil = crop_top_bar(pil, topbar_crop_pct)

    w_img, h_img = pil.size
    if not w_img or not h_img:
        return None
    ar_img = w_img / h_img
    W = int(width)
    H = int(height)
    ar_box = W / H

    if fit == "contain":
        # Aucun crop: on rentre l'image dans la bo√Æte, centr√©e
        if ar_img >= ar_box:
            tgt_w = W
            tgt_h = int(W / ar_img)
        else:
            tgt_h = H
            tgt_w = int(H * ar_img)
        left_off = int(int(left) + (W - tgt_w) / 2)
        top_off = int(int(top) + (H - tgt_h) / 2)
        pic = slide.shapes.add_picture(
            io.BytesIO(_pil_to_bytes(pil, "PNG")),
            Emu(left_off),
            Emu(top_off),
            width=Emu(tgt_w),
            height=Emu(tgt_h),
        )
        return pic

    # "cover": recouvre la bo√Æte (possibles coupes l√©g√®res)
    crop_left = crop_right = crop_top = crop_bottom = 0.0
    if ar_img > ar_box:
        new_w = H * ar_box
        extra = (W - new_w) / W
        crop_left = crop_right = max(0.0, min(0.05, extra / 2))
    elif ar_img < ar_box:
        new_h = W / ar_box
        extra = (H - new_h) / H
        crop_top = crop_bottom = max(0.0, min(0.03, extra / 2))

    pic = slide.shapes.add_picture(
        io.BytesIO(_pil_to_bytes(pil, "PNG")),
        left,
        top,
        width=width,
        height=height,
    )
    try:
        pic.crop_left = float(crop_left)
        pic.crop_right = float(crop_right)
        pic.crop_top = float(crop_top)
        pic.crop_bottom = float(crop_bottom)
    except Exception:
        pass
    return pic


def _iter_shapes(container):
    for shp in getattr(container, "shapes", []):
        yield shp
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shp, "shapes"):
            for child in _iter_shapes(shp):
                yield child


def _shape_text(shp) -> str:
    try:
        if hasattr(shp, "text_frame") and shp.text_frame:
            return (shp.text or "").strip()
    except Exception:
        pass
    return ""


def _is_match_strict(shp, key: str) -> bool:
    k = (key or "").strip().lower()
    if not k:
        return False
    name = (getattr(shp, "name", "") or "").lower().strip()
    alt = (getattr(shp, "alternative_text", "") or "").lower().strip()
    txt = _shape_text(shp).lower().strip()
    return (name == k) or (alt == k) or (txt == k)


def find_placeholder_bbox(slide, key: str):
    for shp in _iter_shapes(slide):
        if _is_match_strict(shp, key):
            return shp, shp.left, shp.top, shp.width, shp.height
    try:
        for shp in _iter_shapes(slide.slide_layout):
            if _is_match_strict(shp, key):
                return None, shp.left, shp.top, shp.width, shp.height
    except Exception:
        pass
    try:
        for shp in _iter_shapes(slide.slide_layout.slide_master):
            if _is_match_strict(shp, key):
                return None, shp.left, shp.top, shp.width, shp.height
    except Exception:
        pass
    return None, None, None, None, None


def detect_placeholders_by_tokens(
    template_bytes: bytes, tokens: List[str]
) -> Dict[str, List[Tuple[int, int, int, int, int]]]:
    prs = Presentation(io.BytesIO(template_bytes))
    outs: Dict[str, List[Tuple[int, int, int, int, int]]] = {
        t.lower(): [] for t in tokens
    }
    for i, slide in enumerate(prs.slides):
        for tok in tokens:
            shp, left, top, width, height = find_placeholder_bbox(slide, tok)
            if left is not None:
                outs[tok.lower()].append((i, left, top, width, height))
    return outs


def fill_template_by_streams_cover(
    template_bytes: bytes,
    streams: dict,
    trim=True,
    topbar_crop_pct=0.0,
) -> bytes:
    """
    Remplit le template par tokens (PH_TBL / PH_LKR_x ...), en utilisant
    pour chaque 'stream' ses propres options (trim/topbar/fit/pre_crop_pct).
    """
    prs = Presentation(io.BytesIO(template_bytes))

    # On d√©tecte toutes les positions de tokens une seule fois
    tokens: List[str] = []
    for s in streams.values():
        tokens.append(s["token_image"])
        if s.get("token_title"):
            tokens.append(s["token_title"])

    pos = detect_placeholders_by_tokens(template_bytes, tokens)

    for s in streams.values():
        t_img = s["token_image"].lower()
        t_tit = (s.get("token_title") or "").lower() or None

        slots = pos.get(t_img, [])
        images = s.get("images") or []
        titles = s.get("titles") or []

        # options sp√©cifiques par stream (avec valeurs par d√©faut globales)
        cfg_trim = bool(s.get("trim", trim))
        cfg_topbar = float(s.get("topbar_crop_pct", topbar_crop_pct))
        cfg_fit = (s.get("fit", "contain") or "contain").lower()
        cfg_prec = s.get("pre_crop_pct", None)  # dict ou None

        n = min(len(images), len(slots))
        for i in range(n):
            si, left, top, width, height = slots[i]
            slide = prs.slides[si]

            # on retire le placeholder si pr√©sent
            kill = None
            for cand in _iter_shapes(slide):
                if _is_match_strict(cand, t_img):
                    kill = cand
                    break
            if kill is not None:
                try:
                    kill._element.getparent().remove(kill._element)
                except Exception:
                    pass

            # On utilise add_picture_fit avec les options par stream
            add_picture_fit(
                slide,
                images[i],
                left,
                top,
                width,
                height,
                trim=cfg_trim,
                topbar_crop_pct=cfg_topbar,
                fit=cfg_fit,
                pre_crop_pct=cfg_prec,
            )

            # Remplissage du titre si demand√©
            if t_tit and i < len(titles) and titles[i]:
                for cand in _iter_shapes(slide):
                    if _is_match_strict(cand, t_tit) and hasattr(cand, "text_frame"):
                        try:
                            cand.text_frame.clear()
                        except Exception:
                            pass
                        cand.text_frame.text = titles[i]
                        break

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def fill_template_sequential_cover(
    template_bytes: bytes,
    images: List[bytes],
    titles: List[str] | None = None,
    image_placeholder_name=PH_IMAGE_MAIN,
    title_placeholder_name=PH_TITLE_MAIN,
    trim=True,
    topbar_crop_pct=0.0,
) -> bytes:
    prs = Presentation(io.BytesIO(template_bytes))
    target_slide_indexes: List[int] = []

    for i, slide in enumerate(prs.slides):
        _, l, t, w, h = find_placeholder_bbox(slide, image_placeholder_name)
        if l is not None:
            target_slide_indexes.append(i)

    if not target_slide_indexes:
        target_slide_indexes = [0]

    while len(target_slide_indexes) < len(images):
        src = prs.slides[target_slide_indexes[-1]]
        new = prs.slides.add_slide(src.slide_layout)
        for shp in list(src.shapes):
            new_el = deepcopy(shp.element)
            new.shapes._spTree.insert_element_before(new_el, "p:extLst")
        target_slide_indexes.append(len(prs.slides) - 1)

    for i, img in enumerate(images):
        slide = prs.slides[target_slide_indexes[i]]
        shp, l, t, w, h = find_placeholder_bbox(slide, image_placeholder_name)

        if l is None:
            slide_w, slide_h = prs.slide_width, prs.slide_height
            l, t, w, h = (
                Inches(0.5),
                Inches(0.6),
                slide_w - Inches(1.0),
                slide_h - Inches(1.2),
            )
        else:
            if shp is not None:
                try:
                    shp._element.getparent().remove(shp._element)
                except Exception:
                    pass

        add_picture_fit(
            slide,
            img,
            l,
            t,
            w,
            h,
            trim=trim,
            topbar_crop_pct=topbar_pct,
            fit=FIT_MODE,
        )

        if titles and i < len(titles) and titles[i]:
            tgt = None
            key = (title_placeholder_name or "").strip().lower()
            for cand in _iter_shapes(slide):
                if _is_match_strict(cand, key):
                    tgt = cand
                    break
            if tgt and hasattr(tgt, "text_frame"):
                try:
                    tgt.text_frame.clear()
                except Exception:
                    pass
                tgt.text_frame.text = titles[i]

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def extract_slide_pictures_from_ppt(ppt_bytes) -> List[bytes]:
    prs = Presentation(io.BytesIO(ppt_bytes))
    imgs: List[bytes] = []
    for slide in prs.slides:
        max_pic = None
        max_area = 0
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                area = int(sh.width) * int(sh.height)
                if area > max_area:
                    max_area = area
                    max_pic = sh
        if max_pic is not None:
            imgs.append(max_pic.image.blob)
    return imgs


def extract_slide_titles_tableau_aware(
    ppt_bytes, workbook_title_guess: Optional[str] = None
) -> List[str]:
    prs = Presentation(io.BytesIO(ppt_bytes))
    titles: List[str] = []
    for slide in prs.slides:
        pic = None
        for sh in slide.shapes:
            if getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                if (pic is None) or (
                    int(sh.width) * int(sh.height)
                    > int(getattr(pic, "width", 0)) * int(getattr(pic, "height", 0))
                ):
                    pic = sh
        pic_top = int(getattr(pic, "top", 10**9)) if pic else None

        candidates = []
        for sh in slide.shapes:
            t = _shape_text(sh)
            if not t:
                continue
            if workbook_title_guess and t.strip() == (workbook_title_guess.strip()):
                continue
            top = int(getattr(sh, "top", 10**9))
            score = abs(top - pic_top) if pic_top is not None else top
            candidates.append((score, -len(t), t))

        titles.append(sorted(candidates)[0][2] if candidates else "")

    return titles


# =====================================
# Gmail / URL ‚Üí PDF
# =====================================
def gmail_service_from_refresh():
    if not GMAIL_AVAILABLE:
        raise RuntimeError("Module Gmail indisponible dans cet environnement.")
    creds = Credentials(
        None,
        refresh_token=GMAIL_REFRESH_TOKEN,
        token_uri="https://oauth2.googleapis.com/token",
        client_id=GMAIL_CLIENT_ID,
        client_secret=GMAIL_CLIENT_SECRET,
        scopes=["https://www.googleapis.com/auth/gmail.readonly"],
    )
    try:
        return build("gmail", "v1", credentials=creds, cache_discovery=False)
    except RefreshError as e:
        logger.error("Gmail OAuth RefreshError : %s", e)
        raise RuntimeError(
            "Gmail OAuth : refresh token invalide ou expir√© (RefreshError). "
            "R√©g√©n√®re le refresh token avec ce client_id / client_secret."
        ) from e


def _iter_all_parts(payload):
    stack = [payload] if payload else []
    while stack:
        p = stack.pop()
        yield p
        for child in (p.get("parts") or []):
            stack.append(child)


def _find_pdf_attachments(msg_payload):
    out = []
    for p in _iter_all_parts(msg_payload):
        filename = (p.get("filename") or "").strip()
        body = p.get("body") or {}
        att_id = body.get("attachmentId")
        mime = (p.get("mimeType") or "").lower()
        is_pdf = filename.lower().endswith(".pdf") or mime == "application/pdf"
        if att_id and is_pdf:
            if not filename:
                filename = "report.pdf"
            out.append((filename, att_id))
    return out


def fetch_latest_looker_pdf_bytes_gmail(query: str):
    try:
        svc = gmail_service_from_refresh()
        user = GMAIL_USER or "me"
        res = (
            svc.users()
            .messages()
            .list(
                userId=user,
                q=query,
                maxResults=20,
                includeSpamTrash=False,
            )
            .execute()
        )
        msgs = res.get("messages", []) or []

        for m in msgs:
            msg = (
                svc.users().messages().get(userId=user, id=m["id"]).execute()
            )
            atts = _find_pdf_attachments(msg.get("payload", {}))
            for (fname, att_id) in atts:
                att = (
                    svc.users()
                    .messages()
                    .attachments()
                    .get(userId=user, messageId=msg["id"], id=att_id)
                    .execute()
                )
                data = base64.urlsafe_b64decode(att["data"])
                logger.info("PDF Looker r√©cup√©r√© via Gmail : %s", fname)
                return data, (fname or "report.pdf")

        logger.warning("Aucun PDF Gmail trouv√© pour la requ√™te : %s", query)
        return None, None

    except RefreshError:
        st.error("Gmail OAuth : refresh token invalide/expir√©.")
        return None, None
    except HttpError as e:
        if getattr(e, "reason", "") == "failedPrecondition":
            st.error("Le compte GMAIL_USER n'a pas de bo√Æte Gmail active.")
            return None, None
        st.error(f"Erreur Gmail API : {e}")
        return None, None
    except Exception as e:
        logger.error("Erreur fetch_latest_looker_pdf_bytes_gmail : %s", e)
        st.error(f"Gmail : {e}")
        return None, None


def _normalize_public_pdf_url(url: str) -> str:
    u = url.strip()
    m = re.search(r"drive\.google\.com/file/d/([^/]+)/", u)
    if m:
        fid = m.group(1)
        return f"https://drive.google.com/uc?export=download&id={fid}"
    q = urllib.parse.urlparse(u).query
    qs = urllib.parse.parse_qs(q)
    if "drive.google.com" in u and "id" in qs:
        fid = qs["id"][0]
        return f"https://drive.google.com/uc?export=download&id={fid}"
    if (".sharepoint.com" in u or "onedrive.live.com" in u) and "download=1" not in u:
        sep = "&" if "?" in u else "?"
        return u + sep + "download=1"
    if "dropbox.com" in u:
        u = re.sub(r"[?&]dl=0\b", "", u)
        if "dl=1" not in u:
            u = u + ("&" if "?" in u else "?") + "dl=1"
        return u
    return u


def fetch_looker_pdf_from_url(url: str, timeout=45):
    try:
        norm = _normalize_public_pdf_url(url)
        r = requests.get(
            norm,
            timeout=timeout,
            allow_redirects=True,
            stream=True,
        )
        r.raise_for_status()
        data = r.content
        if not data.startswith(b"%PDF"):
            raise RuntimeError("Le contenu n'est pas un PDF direct.")
        fname = None
        cd = r.headers.get("content-disposition", "")
        if "filename=" in cd:
            fname = cd.split("filename=")[-1].strip().strip('"')
        if not fname:
            fname = os.path.basename(urllib.parse.urlparse(norm).path) or "report.pdf"
        if not fname.lower().endswith(".pdf"):
            fname += ".pdf"
        logger.info("PDF Looker t√©l√©charg√© via URL : %s", fname)
        return data, fname
    except Exception as e:
        logger.error("Erreur t√©l√©chargement Looker depuis URL : %s", e)
        st.error(f"Erreur t√©l√©chargement depuis URL : {e}")
        return None, None


# =====================================
# APP MAIN
# =====================================
def app_main():
    st.title("Export Pr√©sentation DKAM (Tableau + Looker)")
    st.caption(
        "‚öôÔ∏è Identifiants **Tableau** et sources **Looker** stock√©s localement "
        "(config/local_device_creds.json)."
    )

    # ---------- 1) Identifiants Tableau ----------
    saved = load_local_device_creds() or {}
    st.subheader("üîë Identifiants Tableau")
    server, site, pat_name, pat_secret = tableau_credentials_ui(
        saved, key_prefix="tbl_main"
    )

    colA, colB, colC, colD, colE = st.columns([1, 1, 1, 1, 1])
    if colA.button("üíæ Enregistrer / Mettre √† jour"):
        payload = {
            "server": server,
            "site": site,
            "pat_name": pat_name,
            "pat_secret": pat_secret,
        }
        if not (server and pat_name and pat_secret):
            st.error(
                "Renseigne **server, PAT name et PAT secret** "
                "(site peut √™tre vide si 'Default')."
            )
        else:
            save_local_device_creds(payload)
            logger.info(
                "Identifiants Tableau sauvegard√©s pour server=%s, site=%s, pat_name=%s",
                server,
                site,
                pat_name,
            )
            st.success("Identifiants enregistr√©s localement.")
            st.rerun()
    if colB.button("‚Ü©Ô∏è Recharger depuis le cache local"):
        st.experimental_set_query_params()
        st.rerun()
    if colC.button("üóëÔ∏è Supprimer du cache local"):
        clear_local_device_creds()
        logger.info("Identifiants Tableau supprim√©s du cache local")
        st.success("Identifiants supprim√©s.")
        st.rerun()
    if colD.button("üîÑ Rafra√Æchir workbooks/vues"):
        cached_fetch_workbooks.clear()
        cached_fetch_views.clear()
        logger.info("Cache workbooks/views vid√© manuellement")
        st.success("Cache des listes vid√©.")
        st.experimental_set_query_params()
        st.rerun()
    # colE laiss√© libre pour futures options

    st.markdown("---")

    # ---------- 2) S√©lection Tableau ----------
    creds_ready = all([server, pat_name, pat_secret])
    if not creds_ready:
        st.info(
            "Renseigne/charge tes identifiants (server/PAT) pour afficher "
            "la s√©lection des reportings."
        )
        return

    st.subheader("üìä S√©lection des reportings Tableau")
    with st.expander("üîç Choisir le reporting et la story", expanded=True):
        try:
            with st.spinner("R√©cup√©ration des reportings‚Ä¶"):
                workbooks = cached_fetch_workbooks(
                    server, site, pat_name, pat_secret, TABLEAU_API_VERSION
                )
        except Exception as e:
            logger.error("Erreur r√©cup√©ration workbooks Tableau : %s", e)
            st.error(f"Impossible de r√©cup√©rer les reportings : {e}")
            workbooks = []

        wb_names = [wb["name"] for wb in workbooks]
        selected_wb_name = (
            st.selectbox(
                "Nom du reporting",
                options=wb_names,
                index=0 if wb_names else None,
            )
            if wb_names
            else None
        )
        selected_wb = (
            next((wb for wb in workbooks if wb.get("name") == selected_wb_name), None)
            if selected_wb_name
            else None
        )
        selected_wb_id = selected_wb.get("id") if selected_wb else None

        selected_view = None
        if selected_wb_id:
            try:
                with st.spinner("R√©cup√©ration des stories du reporting‚Ä¶"):
                    views = cached_fetch_views(
                        server,
                        site,
                        pat_name,
                        pat_secret,
                        TABLEAU_API_VERSION,
                        selected_wb_id,
                    )
            except Exception as e:
                logger.error("Erreur r√©cup√©ration vues Tableau : %s", e)
                st.error(f"Impossible de r√©cup√©rer les stories du reporting : {e}")
                views = []

            view_names = [v["name"] for v in views]
            selected_view_name = (
                st.selectbox(
                    "Nom de la story",
                    options=view_names,
                    index=0 if view_names else None,
                )
                if view_names
                else None
            )
            selected_view = (
                next((v for v in views if v.get("name") == selected_view_name), None)
                if selected_view_name
                else None
            )

        if selected_wb and selected_view:
            st.success(
                f"üéØ S√©lection : **{selected_wb['name']}** ‚Üí **{selected_view['name']}**"
            )
            st.caption(
                f"IDs: workbook={selected_wb['id']} | view={selected_view['id']}"
            )

    st.markdown("---")

    # ---------- 3) Template PPT ----------
    st.subheader("üìÑ Template PowerPoint")
    discovered = discover_templates()
    options = ["(aucun)"] + list(discovered.keys())
    tpl_label = st.selectbox("Choisir un template", options=options)
    tpl_bytes, tpl_name = (
        get_template_bytes(tpl_label, discovered)
        if tpl_label != "(aucun)"
        else (None, None)
    )

    if not tpl_bytes:
        st.info("Choisis un template .pptx pour activer la g√©n√©ration.")

    keep_tbl_titles = st.checkbox("Conserver les titres Tableau du template", value=True)
    keep_lkr_titles = st.checkbox("Conserver les titres Looker du template", value=True)

    # Pas de slider : valeur d√©termin√©e par env TOPBAR_CROP_PCT
    topbar_pct = TOPBAR_CROP_PCT

    st.markdown("---")

    # ---------- 4) Looker ‚Äî multi-sources ----------
    with st.expander("üìß Configurer les sources Looker", expanded=True):
        looker_sources_ui()

    # ---------- 5) G√©n√©rer ----------
    c1, _ = st.columns([1, 3])
    disabled = not (selected_wb and selected_view and tpl_bytes)

    if c1.button("üöÄ G√©n√©rer la pr√©sentation", disabled=disabled):
        try:
            logger.info(
                "D√©but g√©n√©ration PPT : wb=%s, view=%s, template=%s",
                selected_wb["name"] if selected_wb else None,
                selected_view["name"] if selected_view else None,
                tpl_name,
            )

            # -- Connect Tableau pour export --
            with st.spinner("Connexion Tableau + export‚Ä¶"):
                sess = TableauSession(
                    server, site, pat_name, pat_secret, TABLEAU_API_VERSION
                )
                sess.signin()

                # On tente d'abord le workbook PPT, sinon la vue PPT
                wb_ppt = sess.export_workbook_ppt(selected_wb["id"])
                if wb_ppt:
                    logger.info(
                        "Export PPT workbook Tableau OK (id=%s)", selected_wb["id"]
                    )
                    all_imgs = extract_slide_pictures_from_ppt(wb_ppt)
                    all_titles = extract_slide_titles_tableau_aware(
                        wb_ppt, workbook_title_guess=selected_wb["name"]
                    )
                    tbl_images = all_imgs
                    tbl_titles_derived = all_titles
                else:
                    logger.info(
                        "Export PPT workbook impossible, tentative sur la vue (view_id=%s)",
                        selected_view["id"],
                    )
                    view_ppt = sess.export_view_ppt(selected_view["id"])
                    if not view_ppt:
                        raise RuntimeError(
                            "Impossible d'obtenir le PPT Tableau (exports d√©sactiv√©s ?)"
                        )
                    tbl_images = extract_slide_pictures_from_ppt(view_ppt)
                    tbl_titles_derived = extract_slide_titles_tableau_aware(
                        view_ppt, workbook_title_guess=selected_wb["name"]
                    )

                sess.signout()

            if not tbl_images:
                logger.warning("Aucune image extraite du PPT Tableau.")
                st.error("Aucune image utilisable depuis le PPT Tableau.")
                return

            tbl_titles = (
                None
                if keep_tbl_titles
                else [
                    (
                        tbl_titles_derived[i]
                        if i < len(tbl_titles_derived)
                        and tbl_titles_derived[i]
                        else f"Slide {i+1}"
                    )
                    for i in range(len(tbl_images))
                ]
            )

            # -- Looker: r√©cup√©rer chaque source (Gmail/URL) -> PDF -> images --
            lkr_all_images: List[List[bytes]] = []
            lkr_all_titles: List[List[str]] = []
            gmail_ready = (
                GMAIL_AVAILABLE
                and GMAIL_CLIENT_ID
                and GMAIL_CLIENT_SECRET
                and GMAIL_REFRESH_TOKEN
            )

            with st.spinner("R√©cup√©ration Looker‚Ä¶"):
                for src in st.session_state.get("lkr_sources", []):
                    mode = src.get("mode")

                    if mode == "gmail":
                        if not gmail_ready:
                            if not st.session_state.get("_gmail_warned"):
                                st.error(
                                    "Gmail OAuth non configur√© correctement "
                                    "(client_id / client_secret / refresh_token)."
                                )
                                st.session_state["_gmail_warned"] = True
                            continue

                        q = build_gmail_query(
                            src.get("senders", ""),
                            src.get("subjects", ""),
                            src.get("label", ""),
                            int(src.get("days", 60)),
                        )
                        pdf_bytes, fname = fetch_latest_looker_pdf_bytes_gmail(q)

                        if pdf_bytes:
                            logger.info(
                                "PDF Looker r√©cup√©r√© via Gmail : %s", fname
                            )
                            imgs = pdf_to_png_bytes(pdf_bytes, zoom=2.0)
                            lkr_all_images.append(imgs)
                            lkr_all_titles.append(
                                [f"Page {i+1}" for i in range(len(imgs))]
                            )
                        else:
                            if not st.session_state.get("_gmail_warned"):
                                st.warning(
                                    f"Aucun PDF via Gmail pour la requ√™te : {q}"
                                )
                                st.session_state["_gmail_warned"] = True

                    elif mode == "url":
                        url_val = (src.get("url") or "").strip()
                        if not url_val:
                            st.warning("URL Looker manquante.")
                            continue

                        pdf_bytes, fname = fetch_looker_pdf_from_url(url_val)
                        if not pdf_bytes:
                            st.warning(
                                "T√©l√©chargement Looker impossible depuis l'URL fournie."
                            )
                            continue

                        logger.info("PDF Looker r√©cup√©r√© via URL : %s", fname)
                        imgs = pdf_to_png_bytes(pdf_bytes, zoom=2.0)
                        lkr_all_images.append(imgs)
                        lkr_all_titles.append(
                            [f"Page {i+1}" for i in range(len(imgs))]
                        )

            # -- Assemblage PPT : bascule placeholders si PH_TBL / PH_LKR_* existent --
            with st.spinner("Assemblage du PowerPoint‚Ä¶"):
                streams: Dict[str, Dict[str, Any]] = {
                    "tbl": {
                        "token_image": PH_TBL,
                        "token_title": PH_TBL_TITLE,
                        "images": tbl_images,
                        "titles": (None if keep_tbl_titles else tbl_titles),
                        # R√©glages Tableau
                        "trim": True,
                        "topbar_crop_pct": topbar_pct,
                        "fit": FIT_MODE,
                        "pre_crop_pct": None,
                    }
                }

                for idx, imgs in enumerate(lkr_all_images, start=1):
                    tok = f"{PH_LKR_BASE}_{idx}"
                    tok_title = f"{PH_LKR_TITLE_BASE}_{idx}"
                    streams[f"lkr_{idx}"] = {
                        "token_image": tok,
                        "token_title": tok_title,
                        "images": imgs,
                        "titles": (
                            None
                            if keep_lkr_titles
                            else lkr_all_titles[idx - 1]
                        ),
                        # R√©glages Looker
                        "trim": LKR_TRIM,
                        "topbar_crop_pct": 0.0,
                        "fit": LKR_FIT_MODE,
                        "pre_crop_pct": {
                            "top": LKR_CROP_TOP,
                            "bottom": LKR_CROP_BOTTOM,
                            "left": LKR_CROP_LEFT,
                            "right": LKR_CROP_RIGHT,
                        },
                    }

                tokens_map = detect_placeholders_by_tokens(
                    tpl_bytes,
                    [PH_TBL, PH_TBL_TITLE, PH_IMAGE_MAIN, PH_TITLE_MAIN]
                    + [
                        f"{PH_LKR_BASE}_{i}"
                        for i in range(1, len(lkr_all_images) + 1)
                    ]
                    + [
                        f"{PH_LKR_TITLE_BASE}_{i}"
                        for i in range(1, len(lkr_all_images) + 1)
                    ],
                )

                has_tbl_slots = bool(tokens_map.get(PH_TBL.lower()))
                has_any_lkr_slots = any(
                    k.startswith(PH_LKR_BASE.lower() + "_")
                    for k in tokens_map.keys()
                )

                if has_tbl_slots or has_any_lkr_slots:
                    # Remplissage par emplacements (PH_TBL / PH_LKR_x)
                    final_ppt = fill_template_by_streams_cover(
                        template_bytes=tpl_bytes,
                        streams=streams,
                        trim=True,
                        topbar_crop_pct=topbar_pct,
                    )
                else:
                    # Fallback s√©quentiel (PH_IMAGE_MAIN / PH_TITLE)
                    flat_lkr = [img for sub in lkr_all_images for img in sub]
                    images = tbl_images + flat_lkr

                    titles = None
                    if not keep_tbl_titles or not keep_lkr_titles:
                        flat_titles = (
                            [t for sub in lkr_all_titles for t in sub]
                            if not keep_lkr_titles
                            else []
                        )
                        titles = (tbl_titles or []) + flat_titles

                    final_ppt = fill_template_sequential_cover(
                        template_bytes=tpl_bytes,
                        images=images,
                        titles=titles,
                        image_placeholder_name=PH_IMAGE_MAIN,
                        title_placeholder_name=PH_TITLE_MAIN,
                        trim=True,
                        topbar_crop_pct=topbar_pct,
                    )

                safe = "".join(
                    c
                    for c in (
                        selected_view.get("contentUrl")
                        or selected_view["name"]
                    )
                    if c.isalnum() or c in "_-"
                )

                logger.info("PPT final g√©n√©r√© avec succ√®s : %s", safe)

                st.success("üéâ Pr√©sentation g√©n√©r√©e !")
                st.download_button(
                    "‚¨áÔ∏è T√©l√©charger la pr√©sentation",
                    data=final_ppt,
                    file_name=f"export_mix_{safe}.pptx",
                    mime=(
                        "application/vnd.openxmlformats-officedocument."
                        "presentationml.presentation"
                    ),
                    use_container_width=True,
                    type="primary",
                )

                if not lkr_all_images:
                    st.warning(
                        "Attention : aucun PDF Looker r√©cup√©r√© ‚Äî la pr√©sentation "
                        "contient uniquement la partie Tableau."
                    )

        except Exception as e:
            logger.error("Erreur lors de la g√©n√©ration PPT : %s", e)
            st.error("Une erreur est survenue pendant la g√©n√©ration.")
            st.code(traceback.format_exc())


# =====================================
# Routing
# =====================================
def guard_and_run():
    if guard_access():
        app_main()


guard_and_run()